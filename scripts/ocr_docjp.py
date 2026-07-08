#!/usr/bin/env python3
"""Extract text from DocJP files (PDF/DOCX/XLSX/XLS/PPTX) into Markdown."""

from __future__ import annotations

import sys
import traceback
from pathlib import Path

SRC = Path("documents/Research/DocJP")
DST = Path("documents/Research/DocJP_md")


def extract_pdf(path: Path) -> str:
    import fitz

    doc = fitz.open(path)
    pages: list[str] = []
    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if text:
            pages.append(f"<!-- page {i} -->\n{text}")
    doc.close()
    return "\n\n".join(pages)


def extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
            rows.insert(1, header_sep)
            parts.append("\n".join(rows))
    return "\n\n".join(parts)


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_data: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(cells):
                rows_data.append(cells)
        if not rows_data:
            continue
        max_cols = max(len(r) for r in rows_data)
        for r in rows_data:
            while len(r) < max_cols:
                r.append("")
        lines: list[str] = [f"## {sheet_name}"]
        for i, row in enumerate(rows_data):
            lines.append("| " + " | ".join(row) + " |")
            if i == 0:
                lines.append("| " + " | ".join(["---"] * max_cols) + " |")
        parts.append("\n".join(lines))
    wb.close()
    return "\n\n".join(parts)


def extract_xls(path: Path) -> str:
    import xlrd

    wb = xlrd.open_workbook(path)
    parts: list[str] = []
    for sheet in wb.sheets():
        rows_data: list[list[str]] = []
        for rx in range(sheet.nrows):
            cells = [str(sheet.cell_value(rx, cx)).strip() for cx in range(sheet.ncols)]
            if any(cells):
                rows_data.append(cells)
        if not rows_data:
            continue
        lines: list[str] = [f"## {sheet.name}"]
        for i, row in enumerate(rows_data):
            lines.append("| " + " | ".join(row) + " |")
            if i == 0:
                lines.append("| " + " | ".join(["---"] * len(row)) + " |")
        parts.append("\n".join(lines))
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                table = shape.table
                rows: list[str] = []
                for row in table.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                    rows.insert(1, header_sep)
                    texts.append("\n".join(rows))
        if texts:
            parts.append(f"<!-- slide {i} -->\n" + "\n\n".join(texts))
    return "\n\n".join(parts)


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".xls": extract_xls,
    ".pptx": extract_pptx,
}


def main() -> None:
    src = Path(sys.argv[1]) if len(sys.argv) > 1 else SRC
    dst = Path(sys.argv[2]) if len(sys.argv) > 2 else DST
    dst.mkdir(parents=True, exist_ok=True)

    files = sorted(src.iterdir())
    ok = 0
    fail = 0
    skip = 0
    empty = 0

    for f in files:
        if f.is_dir():
            continue
        ext = f.suffix.lower()
        extractor = EXTRACTORS.get(ext)
        if extractor is None:
            print(f"  SKIP  {f.name} (unsupported {ext})")
            skip += 1
            continue

        out_name = f.stem + ".md"
        out_path = dst / out_name
        try:
            text = extractor(f)
            if not text.strip():
                print(f"  EMPTY {f.name}")
                out_path.write_text(f"# {f.stem}\n\n(ファイルからテキストを抽出できませんでした)\n", encoding="utf-8")
                empty += 1
            else:
                out_path.write_text(f"# {f.stem}\n\n{text}\n", encoding="utf-8")
                ok += 1
                print(f"  OK    {f.name} -> {out_name} ({len(text)} chars)")
        except Exception:
            print(f"  FAIL  {f.name}")
            traceback.print_exc()
            fail += 1

    print(f"\nDone: {ok} OK, {empty} empty, {fail} failed, {skip} skipped / {len(files)} total")


if __name__ == "__main__":
    main()
