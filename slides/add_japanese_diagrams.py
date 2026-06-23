import base64
import zlib
import urllib.request
from pptx import Presentation
from pptx.util import Inches

def download_mermaid(graph, img_path):
    compressed = zlib.compress(graph.encode('utf-8'), 9)
    b64 = base64.urlsafe_b64encode(compressed).decode('utf-8').rstrip('=')
    url = f"https://kroki.io/mermaid/png/{b64}"
    req = urllib.request.Request(
        url, 
        headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
    )
    with urllib.request.urlopen(req) as response, open(img_path, 'wb') as out_file:
        out_file.write(response.read())

jp_arch = """graph TD
    User(["ユーザー / Client"]) -->|"Webインターフェース"| Frontend["React + Vite Frontend"]
    User -->|"直接呼び出し"| API["FastAPI Backend - Port 8001"]
    Frontend -->|"API呼び出し"| API
    
    subgraph Trai_Tim_He_Thong ["システムのコア (Machine 2)"]
        API -->|"アップロード & クエリ"| RAG["RAG Pipeline"]
        API -->|"コーディングタスク"| Agent["LangGraph Agent"]
        
        RAG -->|"ファイル解析"| Parser["Docling & OCR"]
        RAG -->|"ベクトル生成"| Embedder["BGE-M3 Embedder"]
        RAG -->|"セマンティック検索"| Qdrant[("Qdrant Vector DB - Port 6333")]
        RAG -->|"フォールバック検索"| DuckDuckGo["Web Search ddgs"]
        RAG -->|"MESクエリ"| MES[("MES Database & SQL Agent")]
        
        Agent -->|"利用"| MCP["MCP Client"]
        MCP -->|"ファイルの読み書き"| FS_Server["Filesystem MCP Server"]
        MCP -->|"ソースコード管理"| Git_Server["Git MCP Server"]
        
        RAG -->|"回答の生成"| LiteLLM["LiteLLM Proxy - Port 4000"]
        Agent -->|"推論"| LiteLLM
    end
    
    subgraph Models ["Models (LiteLLMによる調整)"]
        LiteLLM -->|"openai-model / auto"| OpenAI["OpenAI GPT-5.4 mini"]
        LiteLLM -->|"grok-model"| Azure["Azure Grok 4.20 Reasoning"]
        LiteLLM -->|"local-gemma / coding"| Ollama["Machine 1のLocal Gemma4"]
    end
"""

jp_rag = """sequenceDiagram
    participant U as ユーザー
    participant API as FastAPI
    participant Q as Qdrant DB
    participant L as LiteLLM Proxy
    
    U->>API: 質問を送信 (mode: mkac / research)
    API->>API: 質問からベクトルを生成 (BGE-M3)
    API->>Q: セマンティック検索 (ベクトル照合)
    Q-->>API: Top-Kのチャンクを返す (テキスト + 画像URL)
    
    alt Qdrantで見つからない場合
        API->>API: DuckDuckGo Web Search (フォールバック)
    end
    
    API->>API: System Prompt と User Context を構築
    API->>API: 画像のBase64エンコード (存在する場合)
    API->>L: Chat Completionを送信 (Context + 質問)
    L-->>API: Token Streamを返す (Streaming)
    API-->>U: 回答と引用元を徐々に表示
"""

jp_agent = """stateDiagram-v2
    [*] --> Nhan_Yeu_Cau
    Nhan_Yeu_Cau --> Agent_Node : ユーザーがタスクを送信
    
    state "LLM (coding-model)" as Agent_Node {
        direction LR
        Phan_tich --> Lap_Ke_Hoach
        Lap_Ke_Hoach --> Sinh_Lenh_Tool
    }
    
    state "MCP Tools の実行" as Tool_Node
    
    Agent_Node --> Tool_Node : Tool呼び出し (tool_calls)
    Tool_Node --> Agent_Node : 実行結果 (stdout/stderr/file)
    
    Agent_Node --> [*] : Toolの呼び出しを終了し、結果を返す
"""

print("Downloading Japanese Architecture Diagram...")
arch_img = "/home/jkl/Code/VLLM-PD/slides/arch_ja.png"
download_mermaid(jp_arch, arch_img)

print("Downloading Japanese RAG Diagram...")
rag_img = "/home/jkl/Code/VLLM-PD/slides/rag_ja.png"
download_mermaid(jp_rag, rag_img)

print("Downloading Japanese Agent Diagram...")
agent_img = "/home/jkl/Code/VLLM-PD/slides/agent_ja.png"
download_mermaid(jp_agent, agent_img)

print("Updating PPTX...")
ppt_path = "/home/jkl/Code/VLLM-PD/slides/Meibook.pptx"
prs = Presentation(ppt_path)

# Slide Layout
slide_layout = prs.slide_layouts[5]

# Arch Slide
slide_arch = prs.slides.add_slide(slide_layout)
if slide_arch.shapes.title:
    slide_arch.shapes.title.text = "システム全体アーキテクチャ (System Architecture)"
slide_arch.shapes.add_picture(arch_img, Inches(1), Inches(1.5), height=Inches(5.5))

# RAG Slide
slide_rag = prs.slides.add_slide(slide_layout)
if slide_rag.shapes.title:
    slide_rag.shapes.title.text = "RAG データフロー"
slide_rag.shapes.add_picture(rag_img, Inches(1), Inches(1.5), height=Inches(5))

# Agent Slide
slide_agent = prs.slides.add_slide(slide_layout)
if slide_agent.shapes.title:
    slide_agent.shapes.title.text = "Coding Agent のワークフロー"
slide_agent.shapes.add_picture(agent_img, Inches(2), Inches(1.5), height=Inches(5))

prs.save(ppt_path)
print("Successfully added Japanese diagrams to Meibook.pptx")
