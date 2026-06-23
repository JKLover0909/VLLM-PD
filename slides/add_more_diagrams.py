import base64
import zlib
import urllib.request
import sys
from pptx import Presentation
from pptx.util import Inches

def download_mermaid(graph, img_path):
    compressed = zlib.compress(graph.encode('utf-8'), 9)
    b64 = base64.urlsafe_b64encode(compressed).decode('utf-8').rstrip('=')
    url = f"https://kroki.io/mermaid/png/{b64}"
    req = urllib.request.Request(
        url, 
        headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
    )
    with urllib.request.urlopen(req) as response, open(img_path, 'wb') as out_file:
        out_file.write(response.read())

rag_graph = """sequenceDiagram
    participant U as Người dùng
    participant API as FastAPI
    participant Q as Qdrant DB
    participant L as LiteLLM Proxy
    
    U->>API: Gửi câu hỏi (mode: mkac / research)
    API->>API: Tạo Vector từ câu hỏi (BGE-M3)
    API->>Q: Semantic Search (So khớp Vector)
    Q-->>API: Trả về Top-K Chunks (Text + Image URL)
    
    alt Không tìm thấy trong Qdrant
        API->>API: DuckDuckGo Web Search (fallback)
    end
    
    API->>API: Xây dựng System Prompt & User Context
    API->>API: Mã hóa Base64 ảnh (nếu có)
    API->>L: Gửi Chat Completion (Context + Câu hỏi)
    L-->>API: Trả về Token Stream (Streaming)
    API-->>U: Hiển thị dần câu trả lời & Nguồn trích dẫn
"""

agent_graph = """stateDiagram-v2
    [*] --> Nhan_Yeu_Cau
    Nhan_Yeu_Cau --> Agent_Node : Người dùng đưa task
    
    state "LLM (coding-model)" as Agent_Node {
        direction LR
        Phan_tich --> Lap_Ke_Hoach
        Lap_Ke_Hoach --> Sinh_Lenh_Tool
    }
    
    state "Thực thi MCP Tools" as Tool_Node
    
    Agent_Node --> Tool_Node : Gọi tool (tool_calls)
    Tool_Node --> Agent_Node : Kết quả chạy (stdout/stderr/file)
    
    Agent_Node --> [*] : Không gọi tool nữa, trả về kết quả
"""

print("Downloading RAG Diagram...")
rag_img = "/home/jkl/Code/VLLM-PD/slides/rag_diagram.png"
download_mermaid(rag_graph, rag_img)

print("Downloading Agent Diagram...")
agent_img = "/home/jkl/Code/VLLM-PD/slides/agent_diagram.png"
download_mermaid(agent_graph, agent_img)

print("Updating PPTX...")
ppt_path = "/home/jkl/Code/VLLM-PD/slides/Meibook.pptx"
prs = Presentation(ppt_path)

# Add RAG Slide
slide_layout = prs.slide_layouts[5] # Blank / Title only
slide1 = prs.slides.add_slide(slide_layout)
if slide1.shapes.title:
    slide1.shapes.title.text = "Luồng dữ liệu RAG (Retrieval-Augmented Generation)"
slide1.shapes.add_picture(rag_img, Inches(1), Inches(1.5), height=Inches(5))

# Add Agent Slide
slide2 = prs.slides.add_slide(slide_layout)
if slide2.shapes.title:
    slide2.shapes.title.text = "Luồng hoạt động Coding Agent"
slide2.shapes.add_picture(agent_img, Inches(2), Inches(1.5), height=Inches(5))

prs.save(ppt_path)
print("Successfully added all remaining diagrams to Meibook.pptx")
