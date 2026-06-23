import base64
import zlib
import urllib.request
import sys
import subprocess

# 1. Download image from kroki
graph = """graph TD
    User(["Người dùng / Client"]) -->|"Giao diện Web"| Frontend["React + Vite Frontend"]
    User -->|"Gọi trực tiếp"| API["FastAPI Backend - Port 8001"]
    Frontend -->|"Gọi API"| API
    
    subgraph Trai_Tim_He_Thong ["Trái tim Hệ thống (Máy 2)"]
        API -->|"Upload & Query"| RAG["RAG Pipeline"]
        API -->|"Task Coding"| Agent["LangGraph Agent"]
        
        RAG -->|"Phân tích file"| Parser["Docling & OCR"]
        RAG -->|"Tạo Vector"| Embedder["BGE-M3 Embedder"]
        RAG -->|"Truy vấn Semantic"| Qdrant[("Qdrant Vector DB - Port 6333")]
        RAG -->|"Tìm kiếm fallback"| DuckDuckGo["Web Search ddgs"]
        RAG -->|"Truy vấn MES"| MES[("MES Database & SQL Agent")]
        
        Agent -->|"Sử dụng"| MCP["MCP Client"]
        MCP -->|"Đọc/Ghi file"| FS_Server["Filesystem MCP Server"]
        MCP -->|"Quản lý mã nguồn"| Git_Server["Git MCP Server"]
        
        RAG -->|"Tạo câu trả lời"| LiteLLM["LiteLLM Proxy - Port 4000"]
        Agent -->|"Suy luận"| LiteLLM
    end
    
    subgraph Models ["Models (Điều phối bởi LiteLLM)"]
        LiteLLM -->|"openai-model / auto"| OpenAI["OpenAI GPT-5.4 mini"]
        LiteLLM -->|"grok-model"| Azure["Azure Grok 4.20 Reasoning"]
        LiteLLM -->|"local-gemma / coding"| Ollama["Local Gemma4 trên Máy 1"]
    end
"""

compressed = zlib.compress(graph.encode('utf-8'), 9)
b64 = base64.urlsafe_b64encode(compressed).decode('utf-8').rstrip('=')
url = f"https://kroki.io/mermaid/png/{b64}"
print("Downloading from:", url)
img_path = "/home/jkl/Code/VLLM-PD/slides/diagram.png"
req = urllib.request.Request(
    url, 
    headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
)
with urllib.request.urlopen(req) as response, open(img_path, 'wb') as out_file:
    out_file.write(response.read())

# 2. Update PPTX
try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches

ppt_path = "/home/jkl/Code/VLLM-PD/slides/Meibook.pptx"

prs = Presentation(ppt_path)
# Let's add a new slide at the end
slide_layout = prs.slide_layouts[5] # usually a blank slide or title only
slide = prs.slides.add_slide(slide_layout)

# Try to set title if it exists
if slide.shapes.title:
    slide.shapes.title.text = "Kiến trúc hệ thống"

# Insert image
left = Inches(1)
top = Inches(1.5)
height = Inches(5.5)
try:
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
except Exception as e:
    print(f"Failed to add picture: {e}")

prs.save(ppt_path)
print("Successfully added the architecture diagram to Meibook.pptx")
